import os
import base64
import fitz  # PyMuPDF
import docx
import pandas as pd
from pptx import Presentation
from openai import OpenAI
import sys
import logging

# 引用根目录配置
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
import config

# 初始化模块级日志
logger = logging.getLogger(__name__)

KEYS = config.AI_KEYS
CLIENTS = {}
try:
    if KEYS["zhipu"]:
        CLIENTS["zhipu"] = OpenAI(api_key=KEYS["zhipu"], base_url="https://open.bigmodel.cn/api/paas/v4/")
    if KEYS["aliyun"]:
        CLIENTS["aliyun"] = OpenAI(api_key=KEYS["aliyun"], base_url="https://dashscope.aliyuncs.com/compatible-mode/v1")
    if KEYS["deepseek"]:
        CLIENTS["deepseek"] = OpenAI(api_key=KEYS["deepseek"], base_url="https://api.deepseek.com")
    if KEYS["silicon"]:
        CLIENTS["silicon"] = OpenAI(api_key=KEYS["silicon"], base_url="https://api.siliconflow.cn/v1")
except Exception as e:
    logger.warning(f"⚠️ API Client 初始化警告: {e}")

MODELS = {
    "commander": ("deepseek", "deepseek-chat"),   # 主力总结
    "strategist": ("aliyun", "qwen-max"),         # 备用总结
    "hunter": ("zhipu", "glm-4-flash"),           # 快速过滤 (免费/便宜)
    "vision": ("zhipu", "glm-4v-flash")           # 视觉识别
}

class BulletinSummarizer:
    def __init__(self):
        self.clients = CLIENTS
        self.models = MODELS

    def _call_ai(self, role, system_prompt, user_content):
        """通用 AI 调用函数"""
        provider_name, model_name = self.models.get(role, ("deepseek", "deepseek-chat"))
        client = self.clients.get(provider_name)

        if not client:
            logger.warning(f"    ⚠️ 未配置 {provider_name} 的 API Key，跳过 {role}")
            return None

        try:
            temp = config.AI_CONFIG.get("TEMPERATURE", 0.1)
            timeout = config.AI_CONFIG.get("TIMEOUT", 45)
            
            response = client.chat.completions.create(
                model=model_name,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_content}
                ],
                temperature=temp,
                timeout=timeout
            )
            return response.choices[0].message.content
        except Exception as e:
            logger.warning(f"    ⚠️ {role} [{model_name}] 调用失败: {e}")
            return None

    # ==========================
    # 📂 附件解析模块
    # ==========================

    def _extract_pdf(self, filepath):
        text = ""
        try:
            max_pages = config.AI_CONFIG.get("MAX_ATTACH_PAGES", 10)
            with fitz.open(filepath) as doc:
                for page in doc[:max_pages]:
                    text += page.get_text()
            return text[:5000]
        except: return "[PDF解析错误]"

    def _extract_word(self, filepath):
        text = ""
        try:
            doc = docx.Document(filepath)
            for para in doc.paragraphs: text += para.text + "\n"
            return text[:5000]
        except: return "[Word解析错误]"

    def _extract_excel(self, filepath):
        try:
            df = pd.read_excel(filepath, nrows=100).fillna("")
            if df.empty: return "[空Excel表格]"
            return df.to_markdown(index=False)[:4000]
        except Exception as e:
            return f"[Excel解析错误: {str(e)}]"

    def _extract_ppt(self, filepath):
        text = ""
        try:
            max_slides = config.AI_CONFIG.get("MAX_ATTACH_SLIDES", 15)
            prs = Presentation(filepath)
            for slide in prs.slides[:max_slides]:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text[:4000]
        except Exception as e:
            return f"[PPT解析错误: {str(e)}]"

    def _extract_image_content(self, filepath):
        logger.info(f"    👁️ 正在识别图片内容: {os.path.basename(filepath)}...")
        try:
            with open(filepath, "rb") as image_file:
                encoded_string = base64.b64encode(image_file.read()).decode('utf-8')

            client = self.clients.get("zhipu")
            if not client: return "[未配置Vision模型]"

            timeout = config.AI_CONFIG.get("VISION_TIMEOUT", 30)
            response = client.chat.completions.create(
                model="glm-4v-flash",
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": "提取图片中的所有文字，保持原有排版结构。"},
                            {"type": "image_url", "image_url": {"url": encoded_string}}
                        ]
                    }
                ],
                timeout=timeout
            )
            return response.choices[0].message.content
        except Exception as e:
            logger.warning(f"    ⚠️ 图片识别失败: {e}")
            return "[图片无法识别]"

    # ==========================
    # 📉 复杂度优化：原子化处理
    # ==========================

    def _get_extractor_map(self):
        """获取后缀映射表"""
        return {
            '.pdf': self._extract_pdf,
            '.docx': self._extract_word,
            '.doc': self._extract_word,
            '.xlsx': self._extract_excel,
            '.xls': self._extract_excel,
            '.pptx': self._extract_ppt,
            '.ppt': self._extract_ppt,
            '.jpg': self._extract_image_content,
            '.jpeg': self._extract_image_content,
            '.png': self._extract_image_content
        }

    def _process_single_file(self, path, extractors):
        """原子任务：处理单个文件"""
        if not os.path.exists(path):
            return None

        ext = os.path.splitext(path)[1].lower()
        handler = extractors.get(ext)

        if not handler:
            return None

        content = handler(path)
        if not content:
            return None

        return f"\n\n--- 附件 ({os.path.basename(path)}) ---\n{content}\n"

    def process_attachments(self, file_paths):
        """处理附件列表（纯遍历逻辑，复杂度极低）"""
        if not file_paths: return ""

        logger.info(f"    📎 正在预处理 {len(file_paths)} 个附件...")
        extractors = self._get_extractor_map()
        combined_text = ""

        for path in file_paths:
            # 调用原子函数处理单个文件
            result = self._process_single_file(path, extractors)
            if result:
                combined_text += result

        return combined_text

    # ==========================
    # 🧱 原子组件：业务逻辑拆分 (降维打击复杂度)
    # ==========================

    def _build_full_context(self, fetch_result, title):
        """原子任务：组装正文和附件"""
        web_text = fetch_result.get('text', '')
        files = fetch_result.get('files', [])

        # 解析附件
        attach_text = self.process_attachments(files)

        # 确定标题
        safe_title = title if title else (web_text.split('\n')[0] if web_text else "无标题")

        # 组装全文
        full_context = f"【公告标题】: {safe_title}\n\n【网页正文】:\n{web_text}\n{attach_text}"
        return safe_title, full_context

    def _check_relevance(self, safe_title, full_context):
        """原子任务：Hunter 过滤逻辑"""
        # 1. 长度初筛
        if len(full_context) < 20:
            return False

        # 2. 白名单检查
        important_keywords = ["通知", "公告", "公示", "名单", "日程", "安排", "招标", "中标", "竞赛", "讲座", "大创", "补考", "申报"]
        if any(k in safe_title for k in important_keywords):
            logger.info(f"    🛡️ 触发白名单，跳过过滤: {safe_title}")
            return True

        # 3. AI 智能判断
        filter_prompt = """
        你是一个学校通知审核员。请判断以下网页内容是否包含【实质性的通知、新闻、活动或公示信息】。
        
        🔴 判定为 NO (无价值) 的情况：
        1. 仅包含网站导航菜单、页脚、版权声明、友情链接。
        2. 页面提示“404”、“无访问权限”、“系统维护”、“测试页面”。
        3. 正文几乎为空，或仅有“附件”二字但无具体说明。
        4. 纯粹的商业广告。
        
        🟢 判定为 YES (有价值) 的情况：
        1. 包含具体的活动时间、地点、参与人员名单。
        2. 包含科研项目申报、截止日期、招标参数。
        3. 包含具体的新闻报道、会议纪要。
        
        请仅回答 YES 或 NO。
        """
        filter_len = config.AI_CONFIG.get("FILTER_CONTEXT_LEN", 2500)
        is_valuable = self._call_ai("hunter", filter_prompt, full_context[:filter_len])

        if is_valuable and is_valuable.strip().upper().startswith("NO"):
            return False

        return True

    def _generate_summary_content(self, full_context):
        """原子任务：Commander/Strategist 总结逻辑"""
        summary_prompt = """
        你是一个专为高校师生服务的【信息提取助手】。请仔细阅读输入内容，提取关键信息，不要过度概括细节。

        请严格按照以下 Markdown 格式输出：

        📌 **标题**：(原标题，去除非必要修饰)

        🎯 **核心划重点**：
        - (保留具体的研究方向、比赛赛道、招聘岗位等细分列表，不要只写大类！例如：不要只写"农业"，要写"农业(含智慧农业、采后创新等)")
        - (保留具体的硬性要求，如：排名要求、特定专业、必须具备的证书)
        - (保留具体的金额、名额限制)
        - (如果包含"操作指引"，请简述关键步骤，如"需在系统备注栏填写XXX")

        📞 **联系方式**：
        - (提取文中的联系人、电话、邮箱、QQ群、办公地点。如果没有，写"无")

        📎 **附件/链接**：
        - (提取文中出现的重要网址、报名链接、附件名称)

        ⏰ **截止时间**：(精确提取日期和具体时间点)
        """
        max_ctx = config.AI_CONFIG.get("MAX_CONTEXT_LEN", 12000)
        summary = self._call_ai("commander", summary_prompt, full_context[:max_ctx])

        if not summary:
            logger.warning("    ⚠️ Commander 失败，切换 Strategist...")
            summary = self._call_ai("strategist", summary_prompt, full_context[:max_ctx])

        return summary

    # ==========================
    # 🚀 主入口 (重构后结构极简)
    # ==========================

    def summarize(self, fetch_result, title=None):
        if not fetch_result: return None

        # 1. 准备上下文
        safe_title, full_context = self._build_full_context(fetch_result, title)

        # 2. 价值评估 (Hunter)
        if not self._check_relevance(safe_title, full_context):
            return "IGNORE"

        # 3. 生成摘要 (Commander)
        summary = self._generate_summary_content(full_context)

        if not summary:
            return "⚠️ AI 总结失败，请直接查看原文。"

        return summary